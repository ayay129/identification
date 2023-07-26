#! /usr/bin/python3
# -*- coding:utf-8 -*-
# @Time: 2023-06-28 10:37
# @Author: Rangers
# @Site: AIGC生产线，多类型文件识别
# @File: aigc_multi_class.py
from collections import OrderedDict
import os
import platform
import subprocess
import threading
import zipfile
from urllib.parse import urlparse
import rarfile
import requests
from pdf2image import convert_from_bytes
import io
from core.identification import image_procedure
from docx import Document
from config import baidu_client
from openpyxl import load_workbook
from xlrd import open_workbook
from server.bodys import InterfaceError
from pptx import Presentation
from core.pdf_deal import has_images_in_pdf, extract_text_from_pdf


# docx:6390, pdf 4007, jpg 1970,doc 1934,png 1543,
# pptx 520,xlsx 255, txt 98, ppt 59
# xls 59, zip 40, rar 10
# 处理文件，调用接口，

def pdf_to_image_page_stream(data_bytes, page=6):
    # pdf识别内容,使用通用接口,最多识别三也pdf
    if platform.system().lower() == "windows":
        images = convert_from_bytes(data_bytes, poppler_path="D:\\Program Files (x86)\\poppler-23.05.0\\Library\\bin")
    else:
        images = convert_from_bytes(data_bytes)
    if page:
        images = images[:page]
    # 调整大小
    stream = io.BytesIO()
    result_str = ""
    for image in images:
        image.save(stream, format="PNG", quality=90)
        image_stream = stream.getvalue()
        # 改变大小,使之适配
        image_bytes = image_procedure(image_bytes=image_stream)
        resp = baidu_client.accurate(image_bytes)
        if resp.get("error_code"):
            return InterfaceError(code=resp.get("error_code"), message=resp.get("error_msg"))
        results = resp["words_result"]
        if not results:
            continue
        long_strings = [result["words"] for result in results]
        result_str += "\n".join(long_strings) + "\n\n\n"
    return result_str


def image_stream_deal(data_bytes):
    image_bytes = image_procedure(image_bytes=data_bytes)
    resp = baidu_client.accurate(image_bytes)
    if resp.get("error_code"):
        return InterfaceError(code=resp.get("error_code"), message=resp.get("error_msg"))
    results = resp["words_result"]
    if not results:
        return None
    long_strings = [result["words"] for result in results]
    long_strings = "\n".join(long_strings)
    return long_strings


def doc2content(data_bytes):
    # doc文件读取
    thread_id = threading.get_ident()
    doc_file = "{}_temp.doc".format(thread_id)
    with open(doc_file, "wb") as f:
        f.write(data_bytes)
    subprocess.run(["soffice", "--headless", "--convert-to", "docx", doc_file])
    docx_file = "{}_temp.docx".format(thread_id)
    with open(docx_file, "rb") as f:
        docx_binary = f.read()
    content = docx2content(docx_binary)
    os.remove(doc_file)
    os.remove(docx_file)
    return content


def docx2content(data_bytes, page=None):
    document = Document(io.BytesIO(data_bytes))
    if page:
        paragraphs = document.paragraphs[:page]
        tables = document.tables[:page]
    else:
        paragraphs = document.paragraphs
        tables = document.tables
    data = []
    for index, paragraph in enumerate(paragraphs):
        if paragraph.text:
            data.append(paragraph.text)
    result = "\n".join(data)

    table_data = []
    for table in tables:
        # 创建一个二维列表来保存表格数据
        table_rows = []

        # 遍历表格的所有行
        for row in table.rows:
            # 创建一个列表来保存行数据
            row_data = []

            # 遍历行中的所有单元格
            for cell in list(OrderedDict.fromkeys(row.cells)):
                # 将单元格内容添加到行数据列表中
                row_data.append(cell.text)

            # 将行数据添加到表格数据列表中
            table_rows.append("|\t".join(row_data))

        # 将表格数据添加到最终结果中
        table_data.append("\n".join(table_rows))
    response = result + "\n\n\n" + "\n\n".join(table_data)
    return response


def xlsx2content(data_bytes):
    excel_file = io.BytesIO(data_bytes)
    wb = load_workbook(excel_file)
    ws = wb.active
    max_rows = ws.max_row
    max_cols = ws.max_column
    data = []
    for col_index in range(1, max_cols + 1):
        for row_index in range(1, max_rows + 1):
            lattice = ws.cell(row=row_index, column=col_index)
            if not lattice.value:
                continue
            elif not str(lattice.value).strip():
                continue
            data.append(str(lattice.value).strip())
    result_str = "\n".join(data)
    return result_str


def xls2content(data_bytes):
    excel_file = io.BytesIO(data_bytes)
    data = open_workbook(file_contents=excel_file.read())
    table = data.sheets()[0]
    n_rows = table.nrows
    n_cols = table.ncols
    data = []
    for row_num in range(n_rows):
        row = table.row_values(row_num)
        if row:
            for i in range(n_cols):
                if row[i].strip():
                    data.append(row[i].strip())
    result_str = "\n".join(data)
    return result_str


def txt2content(data_bytes):
    stream = io.BytesIO(data_bytes)
    text = stream.read().decode("utf-8")
    return text


def ppt_or_pptx2content(data_bytes):
    stream = io.BytesIO(data_bytes)
    presentation = Presentation(stream)
    data = []
    for slide in presentation.slides:
        if not hasattr(slide, "shapes"):
            continue
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if not text.strip():
                    continue
                data.append(text)
    result_str = "\n".join(data)
    return result_str


# docx:6390, pdf 4007, jpg 1970,doc 1934,png 1543,
# pptx 520,xlsx 255, txt 98, ppt 59
# xls 59, zip 40, rar 10
def zip_uncompress2content(data_bytes):
    memory_file = io.BytesIO(data_bytes)
    data = []
    with zipfile.ZipFile(memory_file, "r") as zip:
        file_list = zip.namelist()
        for file_name in file_list:
            with zip.open(file_name) as file:
                # 读取文件的二进制流
                binary_data = file.read()
                content = judge_url_class(file_name, data_bytes=binary_data)
                if not content:
                    continue
                data.append(content)
    return data


def rar_uncompress2content(data_bytes):
    memory_file = io.BytesIO(data_bytes)
    data = []
    with rarfile.RarFile(memory_file, "r") as rar:
        file_list = rar.namelist()
        for file_name in file_list:
            with rar.open(file_name) as file:
                binary_data = file.read()
                content = judge_url_class(file_name, data_bytes=binary_data)
                if not content:
                    continue
                data.append(content)
    return data


def judge_file_class(url, data_bytes):
    pass


def judge_url_class(url, data_bytes):
    url_path = urlparse(url).path
    if url_path.lower().endswith(".docx"):
        data = docx2content(data_bytes)
    elif url_path.lower().endswith(".doc"):
        data = doc2content(data_bytes)
    elif url_path.lower().endswith(".pdf"):
        if has_images_in_pdf(data_bytes):
            data = pdf_to_image_page_stream(data_bytes)
        else:
            data = extract_text_from_pdf(binary_pdf=data_bytes)
    elif url_path.lower().endswith((".pptx", ".ppt")):
        data = ppt_or_pptx2content(data_bytes)
    elif url_path.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", "heic")):
        data = image_stream_deal(data_bytes)
    elif url_path.lower().endswith(".txt"):
        data = txt2content(data_bytes)
    elif url_path.lower().endswith(".xls"):
        data = xls2content(data_bytes)
    elif url_path.lower().endswith(".xlsx"):
        data = xlsx2content(data_bytes)
    elif url_path.lower().endswith(".zip"):
        data = zip_uncompress2content(data_bytes)
    elif url_path.lower().endswith(".rar"):
        data = rar_uncompress2content(data_bytes)
    else:
        data = None
    return data


def test():
    # resp = requests.get(url)
    # if resp.status_code != 200:
    #     return
    # data = judge_file_class(url, data_bytes=resp.content)
    # for root, dirs, files in os.walk("../data/dismantle"):
    #     for name in files:
    #         file_name = "{}/{}".format(root, name)
    #         print(file_name)
    #         with open(file_name, "rb") as f:
    #             data_bytes = f.read()
    #             data = judge_url_class(file_name, data_bytes)
    #         print(data)
    #         print("-" * 100)
    file_name = "../data/dismantle/英文简历1.pdf"
    with open(file_name, "rb") as f:
        data_bytes = f.read()
    data = judge_url_class(file_name, data_bytes)
    print(data)


def distribute_file_class(url):
    resp = requests.get(url)
    if resp.status_code != 200:
        return
    data = judge_url_class(url, data_bytes=resp.content)
    return data


def read_word_table(data_bytes):
    # 创建一个空列表来保存表格数据
    table_data = []

    # 打开Word文档
    doc = Document(io.BytesIO(data_bytes))

    # 遍历文档中的所有表格
    for table in doc.tables:
        # 创建一个二维列表来保存表格数据
        table_rows = []

        # 遍历表格的所有行
        for row in table.rows:
            # 创建一个列表来保存行数据
            row_data = []

            # 遍历行中的所有单元格
            for cell in row.cells:
                # 将单元格内容添加到行数据列表中
                row_data.append(cell.text)

            # 将行数据添加到表格数据列表中
            table_rows.append(row_data)

        # 将表格数据添加到最终结果中
        table_data.append(table_rows)

    return table_data

# if __name__ == '__main__':
#     test()
#